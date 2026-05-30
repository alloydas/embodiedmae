"""Build the 3-slide .pptx deck for EmbodiedMAE-4M run v3 — 2400-epoch completion.

Mirrors make_slide_pptx_v3.py but pulls metrics from metrics_summary_2400.json
and points at the new slide_figure_v3_2400.png + loss_curves_2400.png.
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


REPO        = Path('/work/mech-ai-scratch/alloy/embodiedmae')
FIG_PATH    = REPO / 'slide_figure_v3_2400.png'
ARCH_PATH   = REPO / 'arch_diagram.png'
CURVES_PATH = REPO / 'outputs/4m_run_v3/loss_curves_2400.png'
AGG_PATH    = REPO / 'slide_dump_v3_2400/metrics_aggregate.json'
SUM_PATH    = REPO / 'outputs/4m_run_v3/metrics_summary_2400.json'
OUT_PATH    = REPO / 'slide_figure_v3_2400.pptx'


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=(0, 0, 0), align=PP_ALIGN.LEFT,
                family=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        if family:
            r.font.name = family
        r.font.color.rgb = RGBColor(*color)
    return tb


def fit_picture(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    lx = left + (max_w - pw) // 2
    ty = top  + (max_h - ph) // 2
    return slide.shapes.add_picture(str(path), lx, ty, pw, ph)


def main():
    agg = json.load(open(AGG_PATH)) if AGG_PATH.exists() else {}
    summary = json.load(open(SUM_PATH))
    best     = summary['best']
    final    = summary['final']
    best_ep  = summary['best_epoch']
    final_ep = summary['final_epoch']

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1 : Problem / Modalities / Architecture ─────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_textbox(s1, Inches(0.4), Inches(0.18), Inches(12.6), Inches(0.55),
                'EmbodiedMAE-4M — Multi-modal MAE for synthetic sorghum  ·  run v3 (completed)',
                size=26, bold=True, color=(20, 40, 90))

    add_textbox(s1, Inches(0.4), Inches(0.80), Inches(12.6), Inches(0.40),
                '(i)  Problem',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s1, Inches(0.6), Inches(1.18), Inches(12.4), Inches(0.95),
                ['Plant phenotyping needs representations that capture geometry, appearance, and growth '
                 'parameters jointly — but real-world labels for all of these together are scarce.',
                 'Pre-train one encoder by masked reconstruction over synthetic sorghum, where every '
                 'sample carries all four modalities aligned, so the encoder learns physically meaningful '
                 'plant structure that downstream phenotyping tasks can reuse.'],
                size=12, color=(40, 40, 40))

    add_textbox(s1, Inches(0.4), Inches(2.20), Inches(12.6), Inches(0.40),
                '(ii)  Modalities',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s1, Inches(0.6), Inches(2.58), Inches(12.4), Inches(1.30),
                ['• RGB — 224×224 rendered image  →  196 patch tokens',
                 '• Depth — 224×224 depth map  →  196 patch tokens (target min-max normalised)',
                 '• Point Cloud — 8 196 pts, unit-sphere normalised  →  196 FPS+kNN tokens',
                 '• Procedural params — 1 plant + up to 24 leaf tokens (stem length, leaf split / length / '
                 'bend angle, …) normalised to [0, 1] from the generator YAML'],
                size=12, color=(40, 40, 40))

    add_textbox(s1, Inches(0.4), Inches(3.95), Inches(12.6), Inches(0.40),
                '(iii)  Architecture — high level',
                size=16, bold=True, color=(20, 40, 90))

    if ARCH_PATH.exists():
        fit_picture(s1, ARCH_PATH, Inches(0.4), Inches(4.35),
                    Inches(12.6), Inches(3.05))
    else:
        add_textbox(s1, Inches(0.6), Inches(4.45), Inches(12.4), Inches(0.6),
                    '(arch_diagram.png missing — regenerate with make_arch_diagram.py)',
                    size=12, color=(160, 0, 0))

    # ── Slide 2 : Reconstruction gallery ──────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, Inches(0.4), Inches(0.10), Inches(12.5), Inches(0.55),
                f'Reconstructions across {agg.get("n_samples", "≥15")} validation plants — best epoch {best_ep}',
                size=22, bold=True, color=(20, 40, 90))

    fit_picture(s2, FIG_PATH, Inches(0.15), Inches(0.70),
                Inches(13.0), Inches(6.7))

    notes2 = s2.notes_slide.notes_text_frame
    notes2.text = (
        f'EmbodiedMAE-4M at the best checkpoint (epoch {best_ep} of the completed '
        f'2400-epoch schedule). Training finished cleanly on 2026-05-28.\n\n'
        f'At eval the Dirichlet allocation was driven with mask_ratio='
        f'{agg.get("mask_ratio", 0):.2f}; per-modality fractions are sampled per '
        f'batch.\n\n'
        f'Aggregate over {agg.get("n_samples", "?")} validation plants:\n'
        f'  • RGB MSE  : {agg.get("rgb_mse", 0):.4f}\n'
        f'  • Depth MSE: {agg.get("depth_mse", 0):.4f}\n'
        f'  • PC Chamfer: {agg.get("pc_chamfer", 0):.6f}\n'
        f'  • Param MAE (all)   : {agg.get("param_mae_all", 0):.4f}\n'
        f'  • Param MAE (masked): {agg.get("param_mae_masked", 0):.6f}\n\n'
        f'Best held-out validation during training:\n'
        f'  • val_loss          : {best["val_loss"]:.4f}\n'
        f'  • RGB MSE           : {best["val_rgb_mse"]:.4f}\n'
        f'  • Depth MSE         : {best["val_depth_mse"]:.4f}\n'
        f'  • PC Chamfer        : {best["val_pc_chamfer"]:.6f}\n'
        f'  • Param MAE (masked): {best["val_param_mae_masked"]:.6f}\n'
        f'  • Param acc@0.05    : {best["val_param_acc05"]*100:.1f}%\n'
    )

    # ── Slide 3 : Training curves + metric table ──────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_textbox(s3, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
                f'Training metrics — run v3 (epochs 1 – {final_ep}, completed)',
                size=24, bold=True, color=(20, 40, 90))

    if CURVES_PATH.exists():
        fit_picture(s3, CURVES_PATH, Inches(0.20), Inches(0.85),
                    Inches(8.0), Inches(6.4))
    else:
        add_textbox(s3, Inches(0.20), Inches(0.85), Inches(8.0), Inches(0.6),
                    '(loss_curves_2400.png missing — run regen_history_v3.py)',
                    size=12, color=(160, 0, 0))

    table_left = Inches(8.50)
    add_textbox(s3, table_left, Inches(0.85), Inches(4.6), Inches(0.45),
                f'Best @ epoch {best_ep}   ·   Final @ epoch {final_ep}',
                size=14, bold=True, color=(20, 40, 90))
    rows = [
        ('Metric',                            'Best        Final'),
        ('val_loss (weighted total)',         f'{best["val_loss"]:.4f}      {final["val_loss"]:.4f}'),
        ('val_rgb (× per-patch MSE)',         f'{best["val_rgb"]:.4f}      {final["val_rgb"]:.4f}'),
        ('val_depth (× per-patch MSE)',       f'{best["val_depth"]:.4f}      {final["val_depth"]:.4f}'),
        ('val_pc (× pc_loss_weight=10)',      f'{best["val_pc"]:.4f}      {final["val_pc"]:.4f}'),
        ('val_text (× spline_w=5)',           f'{best["val_text"]:.4f}      {final["val_text"]:.4f}'),
        ('—',                                 '—'),
        ('RGB MSE  (unweighted)',             f'{best["val_rgb_mse"]:.4f}      {final["val_rgb_mse"]:.4f}'),
        ('Depth MSE (on minmax target)',      f'{best["val_depth_mse"]:.4f}      {final["val_depth_mse"]:.4f}'),
        ('PC Chamfer (bidirectional)',        f'{best["val_pc_chamfer"]:.6f}  {final["val_pc_chamfer"]:.6f}'),
        ('Param MSE',                         f'{best["val_param_mse"]:.4f}      {final["val_param_mse"]:.4f}'),
        ('Param MAE (all leaves)',            f'{best["val_param_mae"]:.4f}      {final["val_param_mae"]:.4f}'),
        ('Param MAE (masked tokens)',         f'{best["val_param_mae_masked"]:.6f}  {final["val_param_mae_masked"]:.6f}'),
        ('Param acc @ |err| < 0.05',          f'{best["val_param_acc05"]*100:.1f} %     {final["val_param_acc05"]*100:.1f} %'),
        ('—',                                 '—'),
        ('Run state',                         'COMPLETED ✓'),
        ('Total epochs',                      f'{final_ep}'),
    ]
    top_y = Inches(1.40)
    line_h = Inches(0.30)
    for label, val in rows:
        add_textbox(s3, table_left,             top_y, Inches(2.9), line_h,
                    label, size=11, color=(60, 60, 60), family='Consolas')
        add_textbox(s3, table_left + Inches(2.9), top_y, Inches(1.7), line_h,
                    val, size=11, bold=True, color=(20, 40, 90), family='Consolas')
        top_y += line_h

    notes3 = s3.notes_slide.notes_text_frame
    notes3.text = (
        f'Run v3 (wandb d0lfi27q) completed all {final_ep} epochs on 2026-05-28 23:01. '
        f'Best val loss {best["val_loss"]:.4f} was recorded at epoch {best_ep}, '
        f'with param acc@0.05 reaching {best["val_param_acc05"]*100:.0f}% — i.e. the '
        f'model effectively memorised the procedural parametrisation from the other '
        f'three modalities at the training mask_ratio of 0.15. PC chamfer dropped to '
        f'{best["val_pc_chamfer"]:.6f}; depth and RGB also continued to improve through '
        f'the cosine annealing tail.\n\n'
        f'Note: this 2400-epoch trajectory is the H200×4 resumed run (10727040). A '
        f'concurrent single-A100 duplicate (10727184) was racing on the same output '
        f'directory and has been cancelled.'
    )

    prs.save(OUT_PATH)
    print(f'Wrote {OUT_PATH}  ({OUT_PATH.stat().st_size/1e6:.2f} MB)')


if __name__ == '__main__':
    main()
