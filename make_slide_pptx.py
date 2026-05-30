"""Build a .pptx deck.

- Slide 1: problem definition + modalities + architecture diagram
- Slide 2: the figure (full bleed) with caption + speaker notes
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


REPO = Path('/work/mech-ai-scratch/alloy/embodiedmae')
FIG_PATH  = REPO / 'slide_figure.png'
ARCH_PATH = REPO / 'arch_diagram.png'
OUT_PATH  = REPO / 'slide_figure.pptx'


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=(0, 0, 0), align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)
    return tb


def main():
    th = json.load(open(REPO / 'outputs/4m_run_v2_again/training_history.json'))
    param_acc05    = max(th['val_param_acc05'])
    param_mae_mskd = min(th['val_param_mae_masked'])
    pc_chamfer     = min(th['val_pc_chamfer'])
    depth_mse      = min(th['val_depth_mse'])
    n_epochs       = len(th['train_loss'])

    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    blank = prs.slide_layouts[6]   # blank layout

    # ── Slide 1 : Problem / Modalities / Architecture diagram ──────────────
    s1 = prs.slides.add_slide(blank)

    # Title
    add_textbox(s1, Inches(0.4), Inches(0.18), Inches(12.6), Inches(0.55),
                'EmbodiedMAE-4M — Multi-modal MAE for synthetic sorghum',
                size=26, bold=True, color=(20, 40, 90))

    # (i) Problem definition
    add_textbox(s1, Inches(0.4), Inches(0.80), Inches(12.6), Inches(0.40),
                '(i)  Problem definition',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s1, Inches(0.6), Inches(1.18), Inches(12.4), Inches(0.95),
                ['Plant phenotyping needs representations that capture geometry, appearance, and growth '
                 'parameters jointly — but real-world labels for all of these together are scarce.',
                 'We pre-train a single encoder by masked-reconstruction over synthetic sorghum, where every '
                 'sample comes with all four modalities aligned, so the encoder learns physically meaningful '
                 'plant structure that downstream phenotyping tasks can reuse.'],
                size=12, color=(40, 40, 40))

    # (ii) The modalities
    add_textbox(s1, Inches(0.4), Inches(2.20), Inches(12.6), Inches(0.40),
                '(ii)  The four modalities',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s1, Inches(0.6), Inches(2.58), Inches(12.4), Inches(1.30),
                ['• RGB — 224×224 rendered image  →  196 patch tokens',
                 '• Depth — 224×224 depth map  →  196 patch tokens (target min-max normalised)',
                 '• Point Cloud — 8,196 pts, unit-sphere normalised  →  196 FPS+kNN tokens',
                 '• Procedural params — 1 plant + up to 24 leaf tokens (stem length, leaf split / length / '
                 'bend angle, …) normalised to [0, 1] from the generator YAML'],
                size=12, color=(40, 40, 40))

    # (iii) Architecture diagram (image)
    add_textbox(s1, Inches(0.4), Inches(3.95), Inches(12.6), Inches(0.40),
                '(iii)  Model architecture — high level',
                size=16, bold=True, color=(20, 40, 90))

    with Image.open(ARCH_PATH) as im:
        aw, ah = im.size
    avail_w = Inches(12.6)
    avail_h = Inches(3.05)
    scale = min(avail_w / aw, avail_h / ah)
    pic_w = int(aw * scale)
    pic_h = int(ah * scale)
    left = int((SW - pic_w) / 2)
    top  = Inches(4.35)
    s1.shapes.add_picture(str(ARCH_PATH), left, top, pic_w, pic_h)

    # ── Slide 2 : the figure ───────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6),
                'Reconstructions across four modalities',
                size=24, bold=True, color=(20, 40, 90))

    # Fit the figure to the slide while preserving aspect ratio
    with Image.open(FIG_PATH) as im:
        iw, ih = im.size
    avail_w = Inches(12.6)
    avail_h = Inches(6.7)
    scale = min(avail_w / iw, avail_h / ih)
    pic_w = int(iw * scale)
    pic_h = int(ih * scale)
    left = int((SW - pic_w) / 2)
    top  = Inches(0.75)
    s2.shapes.add_picture(str(FIG_PATH), left, top, pic_w, pic_h)

    # Speaker notes
    notes = s2.notes_slide.notes_text_frame
    notes.text = (
        f'EmbodiedMAE-4M at checkpoint epoch 760 (best validation loss so far), '
        f'within an ongoing run trained to epoch {n_epochs}/2400 '
        f'({n_epochs/2400*100:.0f}% complete — see the progress bar in the figure).\n\n'
        f'For sample Sorghum_0_06, the Dirichlet allocation masked:\n'
        f'  • 46% of RGB patches\n  • 39% of depth patches\n'
        f'  • 25% of point-cloud tokens (50% of points)\n'
        f'  • Plant procedural parameters (entirely masked)\n\n'
        f'Point-cloud colour is held constant across input/reconstruction/GT so the '
        f'audience reads them as the same object.\n\n'
        f'The point of this slide is the bottom row — the plant parameters '
        f'were not shown to the model, yet it recovers stem length (1.05 vs 1.01), '
        f'and the three position components (ps_x/y/z) within 0.005 of GT — '
        f'using only RGB+depth+PC as evidence.\n\n'
        f'Best held-out val metrics over {n_epochs} epochs:\n'
        f'  • Param accuracy @ |err|<0.05: {param_acc05*100:.1f}%\n'
        f'  • Param MAE (masked, norm):   {param_mae_mskd:.4f}\n'
        f'  • PC Chamfer:                 {pc_chamfer:.4f}\n'
        f'  • Depth MSE (norm):           {depth_mse:.4f}'
    )

    prs.save(OUT_PATH)
    print(f'Wrote {OUT_PATH}  ({OUT_PATH.stat().st_size/1e6:.2f} MB)')


if __name__ == '__main__':
    main()
