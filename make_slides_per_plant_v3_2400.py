"""Build one slide per plant for the EmbodiedMAE-4M v3 reconstruction dump.

For every sample in `slide_dump_v3/sample_*` produce:

  slide_per_plant_v3/<name>.png   — single-plant figure with 4 modality rows
                                     × 3 stage columns + procedural-params panel

and assemble them all into:

  slide_figure_per_plant_v3.pptx   — title slide + 15 plant slides

Layout per plant (landscape, 16:9):

                    Input (masked)    Reconstruction    Ground truth
    RGB             ███               ▒▒▒              ███
    Depth           ███               ▒▒▒              ███
    Point cloud     3D                3D               3D
    Procedural      key:val table     key:val table    key:val table
    params          [MASKED highlighted]  (pred)        (target)
"""

import argparse
import json
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import gridspec
from matplotlib.patches import Rectangle
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ── helpers (cloned/adapted from make_slide_figure.py) ───────────────────────

def load_sample(sample_dir: Path):
    s = sample_dir
    rgb_in   = np.array(Image.open(s / 'inputs' / 'rgb.png'))
    rgb_pred = np.array(Image.open(s / 'outputs' / 'rgb.png'))
    depth_in   = np.load(s / 'inputs' / 'depth.npy')
    depth_pred = np.load(s / 'outputs' / 'depth.npy')
    pc_in   = np.load(s / 'inputs' / 'pointcloud.npy')
    pc_pred = np.load(s / 'outputs' / 'pointcloud.npy')
    m_rgb   = np.load(s / 'masks' / 'rgb_mask.npy').astype(bool)
    m_depth = np.load(s / 'masks' / 'depth_mask.npy').astype(bool)
    m_pc    = np.load(s / 'masks' / 'pc_mask.npy').astype(bool)
    m_text  = np.load(s / 'masks' / 'text_mask.npy').astype(bool)
    with open(s / 'outputs' / 'params.json') as f:
        params = json.load(f)
    with open(s / 'metrics.json') as f:
        metrics = json.load(f)
    return dict(
        name=s.name.replace('sample_', ''),
        rgb_in=rgb_in, rgb_pred=rgb_pred,
        depth_in=depth_in, depth_pred=depth_pred,
        pc_in=pc_in, pc_pred=pc_pred,
        m_rgb=m_rgb, m_depth=m_depth, m_pc=m_pc, m_text=m_text,
        params=params, metrics=metrics,
    )


def apply_patch_mask(img, mask, patch=16, fill=0):
    out = img.copy()
    h, w = out.shape[:2]
    nh, nw = h // patch, w // patch
    m2 = mask.reshape(nh, nw)
    for i in range(nh):
        for j in range(nw):
            if m2[i, j]:
                out[i*patch:(i+1)*patch, j*patch:(j+1)*patch] = fill
    return out


def norm01(arr):
    arr = arr.squeeze()
    lo, hi = float(arr.min()), float(arr.max())
    return (arr - lo) / (hi - lo + 1e-8)


def fps_numpy(xyz, npoint, seed=0):
    rng = np.random.default_rng(seed)
    N = xyz.shape[0]
    cent = np.empty(npoint, dtype=np.int64)
    dist = np.full(N, np.inf)
    far = int(rng.integers(0, N))
    for i in range(npoint):
        cent[i] = far
        d = np.sum((xyz - xyz[far]) ** 2, axis=1)
        dist = np.minimum(dist, d)
        far = int(np.argmax(dist))
    return cent


def knn_assign(xyz, cents_idx, k=32):
    cents = xyz[cents_idx]
    d = np.sum((xyz[None] - cents[:, None]) ** 2, axis=-1)
    return np.argpartition(d, k, axis=1)[:, :k]


def pc_visible_indices(xyz, m_pc, num_tokens=196, group_size=32):
    cents_idx = fps_numpy(xyz, num_tokens, seed=0)
    knn = knn_assign(xyz, cents_idx, group_size)
    pts = set()
    for ti, masked in enumerate(m_pc):
        if not masked:
            pts.update(knn[ti].tolist())
    return np.array(sorted(pts), dtype=np.int64)


def plot_pc(ax, pts, color='seagreen', s=0.55, alpha=0.7,
            elev=18, azim=-75, lim=None):
    ax.scatter(pts[:, 0], pts[:, 2], pts[:, 1], c=color, s=s, alpha=alpha,
               depthshade=False, edgecolors='none')
    ax.set_xticks([]); ax.set_yticks([]); ax.set_zticks([])
    ax.view_init(elev=elev, azim=azim)
    if lim is None:
        lim = max(np.abs(pts).max(), 1e-3) * 1.05
    ax.set_xlim(-lim, lim); ax.set_ylim(-lim, lim); ax.set_zlim(-lim, lim)
    try:
        ax.set_box_aspect((1, 1, 1))
    except Exception:
        pass
    for pane in (ax.xaxis.pane, ax.yaxis.pane, ax.zaxis.pane):
        pane.set_visible(False)
    ax.grid(False)


# ── procedural-params formatting ─────────────────────────────────────────────

_PLANT_SHOW = [
    ('sl',   'sl  (stem length)',  '{:.3f}', ''),
    ('ps_x', 'ps_x',                '{:.3f}', ''),
    ('ps_y', 'ps_y',                '{:.3f}', ''),
    ('ps_z', 'ps_z',                '{:.3f}', ''),
    ('pa',   'pa  (plant angle)',  '{:.2f}', '°'),
    ('pr',   'pr  (plant rot)',    '{:.3f}', ''),
]
_LEAF_SHOW = [
    ('sp', 'sp', '{:.3f}', ''),       # split position
    ('ln', 'ln', '{:.3f}', ''),       # length
    ('ba', 'ba', '{:.1f}', '°'),      # bend angle
    ('ra', 'ra', '{:.1f}', '°'),      # rotation angle
    ('wf', 'wf', '{:.3f}', ''),       # width factor
]


def _row(label, raw_dict, key, fmt, unit):
    val = raw_dict.get(key, 0)
    return (label, fmt.format(val) + unit)


def format_params_table(params, mode, max_leaves=6, leaf_keys=None):
    """mode: 'visible' | 'pred' | 'gt'. Returns list of (label, str).

    `leaf_keys` lets callers shrink the per-leaf field list for compact panels.
    """
    rows = []
    plant = params.get('plant') or {}
    leaf_show = _LEAF_SHOW
    if leaf_keys is not None:
        leaf_show = [(k, lbl, fmt, unit)
                     for (k, lbl, fmt, unit) in _LEAF_SHOW if k in leaf_keys]

    if plant:
        if mode == 'visible' and plant.get('masked'):
            rows.append(('Plant', '[MASKED]'))
            for _ in range(5):
                rows.append(('', ''))
        else:
            src = plant['pred_raw'] if mode == 'pred' else plant['target_raw']
            if src:
                for k, lbl, fmt, unit in _PLANT_SHOW:
                    rows.append(_row(lbl, src, k, fmt, unit))
            else:
                rows.append(('Plant', '(no plant token)'))

    leaves = params.get('leaves', [])[:max_leaves]
    for lf in leaves:
        idx = lf['leaf_index']
        if mode == 'visible' and lf.get('masked'):
            rows.append((f'Leaf {idx}', '[MASKED]'))
            continue
        src = lf['pred_raw'] if mode == 'pred' else lf['target_raw']
        rows.append((f'Leaf {idx}', ''))
        for k, lbl, fmt, unit in leaf_show:
            rows.append((f'  {lbl}', fmt.format(src.get(k, 0)) + unit))
    return rows


def draw_param_panel(ax, rows, mode):
    ax.axis('off')
    y = 0.96
    line_h = 0.072
    for label, val in rows:
        is_masked = (val == '[MASKED]')
        col = 'crimson' if is_masked else ('#1d3a14' if mode == 'pred' else '#0b3954')
        weight = 'bold' if is_masked or (mode == 'pred' and val and not label.startswith(' ')) else 'normal'
        if label:
            ax.text(0.02, y, label, transform=ax.transAxes, fontsize=12,
                    family='monospace', color='black',
                    weight='bold' if not label.startswith(' ') else 'normal')
        if val:
            ax.text(0.58, y, val, transform=ax.transAxes, fontsize=12,
                    family='monospace', color=col, weight=weight)
        y -= line_h
        if y < 0.02:
            break


# ── one plant figure ────────────────────────────────────────────────────────

def render_plant(d, out_path: Path, *,
                 best_epoch=None, max_leaves=6, dpi=160):
    if best_epoch is None:
        best_epoch = BEST['epoch']

    name = d['name']
    rgb_in_masked   = apply_patch_mask(d['rgb_in'], d['m_rgb'], patch=16, fill=0)
    depth_in_n      = norm01(d['depth_in'])
    depth_pred_n    = norm01(d['depth_pred'])
    depth_in_masked = apply_patch_mask(
        (depth_in_n * 255).astype(np.uint8), d['m_depth'], patch=16, fill=0)

    vis_idx = pc_visible_indices(d['pc_in'], d['m_pc'])
    pc_vis  = d['pc_in'][vis_idx]
    pc_lim  = max(np.abs(d['pc_in']).max(), np.abs(d['pc_pred']).max()) * 1.05

    # Bottom row params: 2 leaves × 3 fields keeps the 3 small panels legible.
    # The right-hand side panel handles the full leaf-by-leaf table.
    bot_leaf_keys = ('sp', 'ln', 'ba')
    rows_in   = format_params_table(d['params'], 'visible', max_leaves=2, leaf_keys=bot_leaf_keys)
    rows_pred = format_params_table(d['params'], 'pred',    max_leaves=2, leaf_keys=bot_leaf_keys)
    rows_gt   = format_params_table(d['params'], 'gt',      max_leaves=2, leaf_keys=bot_leaf_keys)

    # ── figure ───────────────────────────────────────────────────────────────
    fig = plt.figure(figsize=(20, 12.5))
    gs = gridspec.GridSpec(
        nrows=5, ncols=5,
        height_ratios=[0.45, 1.0, 1.0, 1.0, 1.7],   # taller params row
        width_ratios=[0.20, 1.0, 1.0, 1.0, 1.25],
        hspace=0.14, wspace=0.07,
        left=0.04, right=0.985, top=0.965, bottom=0.025,
    )

    # ── title ────────────────────────────────────────────────────────────────
    tax = fig.add_subplot(gs[0, :])
    tax.axis('off')
    tax.text(0.5, 0.78,
             f'EmbodiedMAE-4M  ·  plant {name}  ·  best epoch {best_epoch} (val {BEST["val_loss"]:.4f})',
             ha='center', va='center', fontsize=22, weight='bold',
             transform=tax.transAxes)

    m = d['metrics']
    sub = (
        f'mask: RGB {d["m_rgb"].mean()*100:.0f}%   '
        f'Depth {d["m_depth"].mean()*100:.0f}%   '
        f'PC {d["m_pc"].mean()*100:.0f}%   '
        f'text {d["m_text"].mean()*100:.0f}%        '
        f'metrics: RGB MSE={m["rgb_mse"]:.4f}   '
        f'Depth MSE={m["depth_mse"]:.4f}   '
        f'PC Chamfer={m["pc_chamfer"]:.4f}   '
        f'Param MAE (masked)={m["param_mae_masked"]:.4f}'
    )
    tax.text(0.5, 0.32, sub, ha='center', va='center',
             fontsize=11, color='dimgray', family='monospace',
             transform=tax.transAxes)

    # ── column headers ───────────────────────────────────────────────────────
    stage_labels = ['Input (masked)', 'Reconstruction', 'Ground truth']
    for c, lbl in enumerate(stage_labels):
        ax_h = fig.add_subplot(gs[1, 1 + c])
        ax_h.axis('off')
        ax_h.text(0.5, 1.15, lbl, ha='center', va='bottom',
                  fontsize=14, weight='bold', transform=ax_h.transAxes,
                  color='#0b3954')
    ax_h = fig.add_subplot(gs[1, 4])
    ax_h.axis('off')
    ax_h.text(0.5, 1.15, 'Procedural params', ha='center', va='bottom',
              fontsize=14, weight='bold', transform=ax_h.transAxes,
              color='#0b3954')

    # ── row labels ───────────────────────────────────────────────────────────
    mod_labels = ['RGB', 'Depth', 'Point Cloud', 'Spline params']
    for r, lbl in enumerate(mod_labels):
        lax = fig.add_subplot(gs[1 + r, 0])
        lax.axis('off')
        lax.text(0.95, 0.5, lbl, ha='right', va='center',
                 fontsize=14, weight='bold', transform=lax.transAxes,
                 rotation=0)

    # ── RGB row ──────────────────────────────────────────────────────────────
    for c, img in enumerate([rgb_in_masked, d['rgb_pred'], d['rgb_in']]):
        ax = fig.add_subplot(gs[1, 1 + c])
        ax.imshow(img); ax.set_xticks([]); ax.set_yticks([])

    # ── Depth row ────────────────────────────────────────────────────────────
    for c, (img, cm) in enumerate([
        (depth_in_masked, 'viridis'),
        (depth_pred_n,    'viridis'),
        (depth_in_n,      'viridis'),
    ]):
        ax = fig.add_subplot(gs[2, 1 + c])
        ax.imshow(img, cmap=cm); ax.set_xticks([]); ax.set_yticks([])

    # ── PointCloud row ───────────────────────────────────────────────────────
    pc_color = 'seagreen'
    cells = [(pc_vis, 1.0, 0.8), (d['pc_pred'], 0.5, 0.6), (d['pc_in'], 0.5, 0.6)]
    for c, (pts, sz, al) in enumerate(cells):
        ax = fig.add_subplot(gs[3, 1 + c], projection='3d')
        plot_pc(ax, pts, color=pc_color, s=sz, alpha=al, lim=pc_lim)
        ax.margins(0)

    # ── Spline-params row (small text panels) ────────────────────────────────
    for c, rows in enumerate([rows_in, rows_pred, rows_gt]):
        ax = fig.add_subplot(gs[4, 1 + c])
        mode = ['visible', 'pred', 'gt'][c]
        draw_param_panel(ax, rows, mode)

    # ── Right-column: full leaf-by-leaf params panel ─────────────────────────
    big_ax = fig.add_subplot(gs[1:, 4])
    big_ax.axis('off')

    # Build a side-by-side target / pred / |Δ| table for ALL real leaves
    plant = d['params'].get('plant') or {}
    leaves = d['params'].get('leaves', [])
    lines = []
    lines.append(('Plant', 'target',  'pred', '|Δ|', plant.get('masked', False)))
    if plant:
        for k, lbl, fmt, unit in _PLANT_SHOW:
            tgt = plant['target_raw'].get(k, 0) if plant.get('target_raw') else 0
            prd = plant['pred_raw'].get(k, 0)   if plant.get('pred_raw')   else 0
            lines.append((f'  {lbl}',
                          fmt.format(tgt) + unit,
                          fmt.format(prd) + unit,
                          fmt.format(abs(tgt-prd)) + unit,
                          plant.get('masked', False)))
    lines.append(('', '', '', '', False))
    for lf in leaves[:max_leaves]:
        masked = lf.get('masked', False)
        marker = ' (masked)' if masked else ''
        lines.append((f'Leaf {lf["leaf_index"]}{marker}', 'target', 'pred', '|Δ|', masked))
        for k, lbl, fmt, unit in _LEAF_SHOW:
            tgt = lf['target_raw'].get(k, 0)
            prd = lf['pred_raw'].get(k, 0)
            lines.append((f'  {lbl}',
                          fmt.format(tgt) + unit,
                          fmt.format(prd) + unit,
                          fmt.format(abs(tgt-prd)) + unit,
                          masked))

    # Render lines
    y = 0.99
    line_h = 0.022
    # Column header
    big_ax.text(0.02, y, f'All visible + first {max_leaves} leaves',
                transform=big_ax.transAxes, fontsize=11, weight='bold',
                family='monospace', color='#0b3954')
    y -= line_h * 1.4
    for label, tgt, prd, diff, masked in lines:
        is_header = (tgt == 'target')
        col_label = 'crimson' if masked and is_header else 'black'
        weight = 'bold' if is_header else 'normal'
        big_ax.text(0.02, y, label, transform=big_ax.transAxes,
                    fontsize=9, family='monospace', color=col_label, weight=weight)
        if not is_header and tgt:
            big_ax.text(0.42, y, tgt, transform=big_ax.transAxes,
                        fontsize=9, family='monospace', color='#0b3954')
            big_ax.text(0.62, y, prd, transform=big_ax.transAxes,
                        fontsize=9, family='monospace',
                        color='crimson' if masked else '#1d3a14')
            big_ax.text(0.84, y, diff, transform=big_ax.transAxes,
                        fontsize=9, family='monospace',
                        color='gray')
        elif is_header and tgt:
            big_ax.text(0.42, y, tgt, transform=big_ax.transAxes,
                        fontsize=9, family='monospace', weight='bold', color='dimgray')
            big_ax.text(0.62, y, prd, transform=big_ax.transAxes,
                        fontsize=9, family='monospace', weight='bold', color='dimgray')
            big_ax.text(0.84, y, diff, transform=big_ax.transAxes,
                        fontsize=9, family='monospace', weight='bold', color='dimgray')
        y -= line_h
        if y < 0.01:
            break

    fig.savefig(out_path, dpi=dpi, bbox_inches='tight', facecolor='white')
    plt.close(fig)


# ── PPTX assembly ────────────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=(0, 0, 0), align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)


def fit_picture(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    lx = left + (max_w - pw) // 2
    ty = top  + (max_h - ph) // 2
    slide.shapes.add_picture(str(path), lx, ty, pw, ph)


_SUM_PATH = Path('/work/mech-ai-scratch/alloy/embodiedmae/outputs/4m_run_v3/metrics_summary_2400.json')
_summary  = json.load(open(_SUM_PATH))
_best     = _summary['best']
BEST = dict(
    epoch=_summary['best_epoch'],
    total_epochs=_summary['final_epoch'],
    cur_epoch=_summary['final_epoch'],
    val_loss=_best['val_loss'],     val_rgb=_best['val_rgb'],
    val_depth=_best['val_depth'],   val_pc=_best['val_pc'],
    val_text=_best['val_text'],
    rgb_mse=_best['val_rgb_mse'],   depth_mse=_best['val_depth_mse'],
    pc_chamfer=_best['val_pc_chamfer'], pc_emd=_best['val_pc_emd'],
    param_mse=_best['val_param_mse'], param_mae=_best['val_param_mae'],
    param_mae_masked=_best['val_param_mae_masked'],
    param_acc05=_best['val_param_acc05'],
)
ARCH_PATH   = Path('/work/mech-ai-scratch/alloy/embodiedmae/arch_diagram.png')
CURVES_PATH = Path('/work/mech-ai-scratch/alloy/embodiedmae/outputs/4m_run_v3/loss_curves_2400.png')


def add_arch_slide(prs, blank, n_plants):
    s = prs.slides.add_slide(blank)
    add_textbox(s, Inches(0.4), Inches(0.18), Inches(12.6), Inches(0.55),
                'EmbodiedMAE-4M — Multi-modal MAE for synthetic sorghum  ·  run v3 (completed)',
                size=26, bold=True, color=(20, 40, 90))

    add_textbox(s, Inches(0.4), Inches(0.80), Inches(12.6), Inches(0.40),
                '(i)  Problem',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s, Inches(0.6), Inches(1.18), Inches(12.4), Inches(0.95),
                ['Plant phenotyping needs representations that capture geometry, appearance, and growth '
                 'parameters jointly — but real-world labels for all of these together are scarce.',
                 'Pre-train one encoder by masked reconstruction over synthetic sorghum, where every '
                 'sample carries all four modalities aligned, so the encoder learns physically meaningful '
                 'plant structure that downstream phenotyping tasks can reuse.'],
                size=12, color=(40, 40, 40))

    add_textbox(s, Inches(0.4), Inches(2.20), Inches(12.6), Inches(0.40),
                '(ii)  Modalities',
                size=16, bold=True, color=(20, 40, 90))
    add_textbox(s, Inches(0.6), Inches(2.58), Inches(12.4), Inches(1.30),
                ['• RGB — 224×224 rendered image  →  196 patch tokens',
                 '• Depth — 224×224 depth map  →  196 patch tokens (target min-max normalised)',
                 '• Point Cloud — 8 196 pts, unit-sphere normalised  →  196 FPS+kNN tokens',
                 '• Procedural params — 1 plant + up to 24 leaf tokens (stem length, leaf split / length / '
                 'bend angle, …) normalised to [0, 1] from the generator YAML'],
                size=12, color=(40, 40, 40))

    add_textbox(s, Inches(0.4), Inches(3.95), Inches(12.6), Inches(0.40),
                '(iii)  Architecture — high level',
                size=16, bold=True, color=(20, 40, 90))

    if ARCH_PATH.exists():
        fit_picture(s, ARCH_PATH, Inches(0.4), Inches(4.35),
                    Inches(12.6), Inches(3.05))
    else:
        add_textbox(s, Inches(0.6), Inches(4.45), Inches(12.4), Inches(0.6),
                    '(arch_diagram.png missing — regenerate with make_arch_diagram.py)',
                    size=12, color=(160, 0, 0))


def add_metrics_slide(prs, blank, agg):
    s = prs.slides.add_slide(blank)
    add_textbox(s, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55),
                f'Training metrics — run v3 (epochs 1 – {BEST["total_epochs"]}, completed)',
                size=24, bold=True, color=(20, 40, 90))

    if CURVES_PATH.exists():
        fit_picture(s, CURVES_PATH, Inches(0.20), Inches(0.85),
                    Inches(8.0), Inches(6.4))
    else:
        add_textbox(s, Inches(0.20), Inches(0.85), Inches(8.0), Inches(0.6),
                    '(loss_curves.png missing — regenerate from training_history)',
                    size=12, color=(160, 0, 0))

    table_left = Inches(8.50)
    add_textbox(s, table_left, Inches(0.85), Inches(4.6), Inches(0.45),
                f'Best @ epoch {BEST["epoch"]}',
                size=16, bold=True, color=(20, 40, 90))
    rows = [
        ('val_loss (weighted total)',         f'{BEST["val_loss"]:.4f}'),
        ('val_rgb (× per-patch MSE)',         f'{BEST["val_rgb"]:.4f}'),
        ('val_depth (× per-patch MSE)',       f'{BEST["val_depth"]:.4f}'),
        ('val_pc (× pc_loss_weight=10)',      f'{BEST["val_pc"]:.4f}'),
        ('val_text (× spline_w=5)',           f'{BEST["val_text"]:.4f}'),
        ('—',                                 '—'),
        ('RGB MSE  (unweighted)',             f'{BEST["rgb_mse"]:.4f}'),
        ('Depth MSE (on minmax target)',      f'{BEST["depth_mse"]:.4f}'),
        ('PC Chamfer (bidirectional)',        f'{BEST["pc_chamfer"]:.6f}'),
        ('Param MSE',                         f'{BEST["param_mse"]:.4f}'),
        ('Param MAE (all leaves)',            f'{BEST["param_mae"]:.4f}'),
        ('Param MAE (masked tokens)',         f'{BEST["param_mae_masked"]:.4f}'),
        ('Param acc @ |err| < 0.05',          f'{BEST["param_acc05"]*100:.1f} %'),
        ('—',                                 '—'),
        ('Eval mask_ratio',                   f'{agg.get("mask_ratio", 0):.2f}'),
        ('Eval RGB MSE (15 plants)',          f'{agg.get("rgb_mse", 0):.4f}'),
        ('Eval Depth MSE',                    f'{agg.get("depth_mse", 0):.4f}'),
        ('Eval PC Chamfer',                   f'{agg.get("pc_chamfer", 0):.4f}'),
        ('Eval Param MAE (all)',              f'{agg.get("param_mae_all", 0):.4f}'),
        ('Eval Param MAE (masked)',           f'{agg.get("param_mae_masked", 0):.4f}'),
        ('—',                                 '—'),
        ('Run state',                         'COMPLETED ✓'),
        ('Total epochs',                      f'{BEST["total_epochs"]}'),
    ]
    top_y = Inches(1.40)
    line_h = Inches(0.26)
    for label, val in rows:
        add_textbox(s, table_left,             top_y, Inches(2.9), line_h,
                    label, size=10, color=(60, 60, 60))
        add_textbox(s, table_left + Inches(2.9), top_y, Inches(1.7), line_h,
                    val, size=10, bold=True, color=(20, 40, 90))
        top_y += line_h


def build_pptx(png_paths, samples, out_path, agg=None):
    if agg is None:
        agg = {}
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Architecture / problem / modalities
    add_arch_slide(prs, blank, len(samples))

    # Slide 2: Metrics summary
    add_metrics_slide(prs, blank, agg)

    # Slides 3+: per-plant
    for png, samp in zip(png_paths, samples):
        s = prs.slides.add_slide(blank)
        add_textbox(s, Inches(0.4), Inches(0.10), Inches(12.5), Inches(0.5),
                    f'Plant: {samp["name"]}',
                    size=20, bold=True, color=(20, 40, 90))
        fit_picture(s, png, Inches(0.15), Inches(0.65),
                    Inches(13.0), Inches(6.75))

        notes = s.notes_slide.notes_text_frame
        m = samp['metrics']
        notes.text = (
            f'Plant {samp["name"]}  (run v3, best epoch {BEST["epoch"]})\n\n'
            f'Mask fractions:\n'
            f'  RGB    : {samp["m_rgb"].mean()*100:.1f}%\n'
            f'  Depth  : {samp["m_depth"].mean()*100:.1f}%\n'
            f'  PC     : {samp["m_pc"].mean()*100:.1f}%\n'
            f'  text   : {samp["m_text"].mean()*100:.1f}%\n\n'
            f'Per-sample metrics:\n'
            f'  RGB MSE          : {m["rgb_mse"]:.4f}\n'
            f'  Depth MSE        : {m["depth_mse"]:.4f}\n'
            f'  PC Chamfer       : {m["pc_chamfer"]:.4f}\n'
            f'  Param MAE (all)  : {m["param_mae_all"]:.4f}\n'
            f'  Param MAE (masked): {m["param_mae_masked"]:.4f}\n'
        )

    prs.save(out_path)
    print(f'Wrote {out_path}  ({out_path.stat().st_size/1e6:.2f} MB)')


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--dump_dir',  default='slide_dump_v3_2400')
    ap.add_argument('--out_dir',   default='slide_per_plant_v3_2400')
    ap.add_argument('--pptx',      default='slide_figure_per_plant_v3_2400.pptx')
    ap.add_argument('--n_plants',  type=int, default=15)
    ap.add_argument('--max_leaves', type=int, default=6)
    ap.add_argument('--dpi',       type=int, default=140)
    args = ap.parse_args()

    dump = Path(args.dump_dir)
    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    sample_dirs = sorted(dump.glob('sample_*'))[:args.n_plants]
    print(f'Found {len(sample_dirs)} samples')

    agg = {}
    agg_path = dump / 'metrics_aggregate.json'
    if agg_path.exists():
        agg = json.load(open(agg_path))

    png_paths = []
    samples   = []
    for sd in sample_dirs:
        d = load_sample(sd)
        png = out_dir / f'{d["name"]}.png'
        print(f'  → {png}')
        render_plant(d, png, max_leaves=args.max_leaves, dpi=args.dpi)
        png_paths.append(png)
        samples.append(d)

    build_pptx(png_paths, samples, Path(args.pptx), agg=agg)


if __name__ == '__main__':
    main()
